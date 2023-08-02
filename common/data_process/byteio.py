# Write the modified data back to the file
'''
with open(file_path, 'wb') as f:
    f.write(modified_data)
with open(file_path, 'rb') as f:
    # Read the data from the file
    data = f.read()
    # Create a BytesIO object from the data
    byte_stream = BytesIO(data)
    modified_data = byte_stream.getvalue()
    print(modified_data)
提取某一个文件所有的文本，然后交给langchain切割文本的函数进行流水线处理
'''
import io
from io import BytesIO
import os
import docx
import langchain
import openai
from typing import List, Any, Optional
import re
import docx2txt
from langchain.docstore.document import Document
import fitz
from hashlib import md5
from markdown import markdown
from abc import abstractmethod, ABC
from copy import deepcopy
import docx2txt
from bs4 import BeautifulSoup
from pptx import Presentation
import pytesseract
from PIL import Image
from paddleocr import PaddleOCR, draw_ocr
import numpy as np
import whisper
from moviepy.editor import VideoFileClip
import time


class File(ABC):
    def __init__(self,
                 id: str,
                 name: Optional[str] = None,
                 meta: Optional[dict[str, Any]] = None,
                 docs: Optional[str] = None):
        self.id = id
        self.name = name or ""
        self.meta = meta or {}
        self.docs = docs or ""  # all raw text

    @classmethod
    @abstractmethod
    def from_bytes(cls, files: BytesIO) -> "File":
        pass

    @staticmethod
    def read_file_to_byte(file_path: str) -> (BytesIO, bool):
        byte_io = None
        try:
            with open(file_path, 'rb') as file:
                byte_io = io.BytesIO(file.read())
            byte_io.seek(0)
            return byte_io, True
        except Exception as e:
            return byte_io, False

    def __repr__(self):
        return f"name:{self.name}, id:{self.id}, docs:{str(self.docs)}"


def strip_consecutive_newlines(text: str) -> str:
    """Strips consecutive newlines from a string
    possibly with whitespace in between
    """
    return re.sub(r"\s*\n\s*", "\n", text)


# package list: docx(complex), docx2txt(simple)
class DocxFile(File):
    @classmethod
    def from_bytes(cls, files: BytesIO) -> "DocxFile":
        text = docx2txt.process(files)
        docs = strip_consecutive_newlines(text)
        return cls(id=md5(files.read()).hexdigest(), docs=docs)


class TxtFile(File):
    @classmethod
    def from_bytes(cls, files: BytesIO) -> "TxtFile":
        files.seek(0)
        text = files.read().decode("utf-8")
        text = strip_consecutive_newlines(text)
        docs = text.strip()  # TODO: consider "\n"
        return cls(id=md5(files.read()).hexdigest(), docs=docs)


class PdfFile(File):
    @classmethod
    def from_bytes(cls, files: BytesIO) -> "PdfFile":
        files.seek(0)
        pdf = fitz.open(stream=files.read(), filetype="pdf")
        docs = " ".join([page.get_text().strip() for page in pdf])
        docs = strip_consecutive_newlines(docs)
        return cls(id=md5(files.read()).hexdigest(), docs=docs)


class MdFile(File):
    @classmethod
    def from_bytes(cls, files: BytesIO) -> "MdFile":
        files.seek(0)
        text = files.read().decode("utf-8")
        html = markdown(text)
        soup = BeautifulSoup(html, 'html.parser')
        docs = soup.get_text('\n')
        docs = strip_consecutive_newlines(docs)
        files.seek(0)
        return cls(id=md5(files.read()).hexdigest(), docs=docs)


class PptFile(File):
    @classmethod
    def from_bytes(cls, files: BytesIO) -> "PptFile":
        files.seek(0)
        prs = Presentation(files)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text += run.text
        text = strip_consecutive_newlines(text)
        files.seek(0)
        return cls(id=md5(files.read()).hexdigest(), docs=text)


class PaddleOCRPdfFile(File):
    @classmethod
    def from_bytes(cls, files: BytesIO) -> "PaddleOCRPdfFile":
        files.seek(0)
        pdf = fitz.open(stream=files.read(), filetype="pdf")
        ocr_model = PaddleOCR(use_gpu=False)  # Set use_gpu=True if you have GPU
        docs = ""
        for page in range(len(pdf)):
            pix = pdf[page].get_pixmap()
            mode = "RGBA" if pix.alpha else "RGB"
            img = Image.frombytes(mode, [pix.width, pix.height], pix.samples)
            img_np = np.array(img)
            ocr_result = ocr_model.ocr(img_np)
            for line in ocr_result:
                line_text = [word_info[-1] for word_info in line]
                text = ''.join([line[0] for line in line_text])
                print(text)
                docs += "\n" + text
        docs = strip_consecutive_newlines(docs)
        return cls(id=md5(files.read()).hexdigest(), docs=docs)


def extract_audio_from_video(video_path, audio_path):
    video_clip = VideoFileClip(video_path)
    audio_clip = video_clip.audio
    audio_clip.write_audiofile(audio_path)
    audio_clip.close()
    video_clip.close()


def extract_content_from_audio(file_path: str, model_size="small"):
    model = whisper.load_model(model_size)
    result = model.transcribe(file_path)
    return result["text"]


class AudioFile():
    def from_path(self, file_path) -> 'AudioFile':
        pass


class VideoFile():
    def from_path(self, file_path) -> 'VideoFile':
        pass
